"""
Presentation Tools Pack

A comprehensive set of tools for creating and editing PowerPoint presentations:
- Create new presentations
- Add slides with various layouts
- Add text, images, tables, and charts
- Apply themes and formatting
- Save presentations
"""

from __future__ import annotations

import json as _json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

import structlog

from squadron.connectivity.mcp_host import mcp_tool

logger = structlog.get_logger(__name__)


@dataclass
class SlideInfo:
    """Information about a created slide."""
    slide_index: int
    layout: str
    title: str | None = None


@dataclass
class PresentationInfo:
    """Information about a presentation."""
    path: str
    slide_count: int
    title: str | None = None


class PresentationTools:
    """
    Presentation Tools Pack.

    Provides tools for creating and editing PowerPoint presentations:
    - Create presentations with custom themes
    - Add slides with various content layouts
    - Add text boxes, images, tables, and charts
    - Format slides with colors and fonts
    - Save and export presentations

    Requirements:
        pip install python-pptx

    Example:
        ```python
        tools = PresentationTools(workspace_root="/path/to/project")

        # Create a new presentation
        info = await tools.create_presentation("my_deck.pptx", title="My Presentation")

        # Add a title slide
        await tools.add_slide("my_deck.pptx", layout="title", title="Welcome", subtitle="Subtitle")

        # Add a content slide
        await tools.add_slide("my_deck.pptx", layout="content", title="Agenda",
                               body="Item 1\\nItem 2\\nItem 3")

        # Save the presentation
        await tools.save_presentation("my_deck.pptx")
        ```
    """

    def __init__(
        self,
        workspace_root: str | Path | None = None,
    ):
        """
        Initialize presentation tools.

        Args:
            workspace_root: Root directory for file operations
        """
        self.workspace_root = Path(workspace_root) if workspace_root else Path.cwd()
        # In-memory store for open presentations: filename -> pptx Presentation object
        self._presentations: dict[str, Any] = {}

    def _resolve_path(self, path: str | Path) -> Path:
        """Resolve a path relative to workspace root."""
        path = Path(path)
        if path.is_absolute():
            return path
        return self.workspace_root / path

    @staticmethod
    def _coerce_int(value: Any, name: str) -> int:
        """Coerce a value to int (LLMs sometimes pass numeric args as strings)."""
        if isinstance(value, int):
            return value
        try:
            return int(value)
        except (TypeError, ValueError):
            raise TypeError(f"Parameter '{name}' must be an integer, got {type(value).__name__}: {value!r}")

    @staticmethod
    def _coerce_float(value: Any, name: str) -> float:
        """Coerce a value to float (LLMs sometimes pass numeric args as strings)."""
        if isinstance(value, (int, float)):
            return float(value)
        try:
            return float(value)
        except (TypeError, ValueError):
            raise TypeError(f"Parameter '{name}' must be a number, got {type(value).__name__}: {value!r}")

    def _get_pptx(self) -> Any:
        """Import python-pptx and return the module."""
        try:
            import pptx
            return pptx
        except ImportError:
            raise ImportError(
                "python-pptx package required. Run: pip install python-pptx"
            )

    def _get_or_load_presentation(self, filename: str) -> Any:
        """Get an in-memory presentation or load from disk."""
        pptx = self._get_pptx()
        if filename not in self._presentations:
            file_path = self._resolve_path(filename)
            if file_path.exists():
                self._presentations[filename] = pptx.Presentation(str(file_path))
            else:
                raise FileNotFoundError(
                    f"Presentation not found: {filename}. "
                    "Use create_presentation() first."
                )
        return self._presentations[filename]

    @mcp_tool(description="Create a new PowerPoint presentation")
    async def create_presentation(
        self,
        filename: str,
        title: str | None = None,
        theme: str = "default",
        width_inches: float = 13.33,
        height_inches: float = 7.5,
    ) -> PresentationInfo:
        """
        Create a new PowerPoint presentation.

        Args:
            filename: Output filename (e.g. 'slides.pptx')
            title: Optional presentation title (added as first slide)
            theme: Theme name ('default' uses built-in blank theme)
            width_inches: Slide width in inches (default: 13.33 widescreen)
            height_inches: Slide height in inches (default: 7.5 widescreen)

        Returns:
            PresentationInfo with path and slide count
        """
        pptx = self._get_pptx()
        from pptx.util import Inches

        # Coerce numeric parameters (LLMs may pass them as strings)
        width_inches = self._coerce_float(width_inches, "width_inches")
        height_inches = self._coerce_float(height_inches, "height_inches")

        prs = pptx.Presentation()
        prs.slide_width = Inches(width_inches)
        prs.slide_height = Inches(height_inches)

        self._presentations[filename] = prs

        slide_count = 0
        if title:
            # Add a title slide
            slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            slide_count = 1

        # Save immediately
        file_path = self._resolve_path(filename)
        file_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(file_path))

        logger.info("Presentation created", filename=filename, slides=slide_count)
        return PresentationInfo(path=str(file_path), slide_count=slide_count, title=title)

    @mcp_tool(description="Add a slide to a PowerPoint presentation")
    async def add_slide(
        self,
        filename: str,
        layout: str = "content",
        title: str | None = None,
        body: str | None = None,
        subtitle: str | None = None,
    ) -> SlideInfo:
        """
        Add a slide to an existing presentation.

        Args:
            filename: Presentation filename
            layout: Slide layout name. Options:
                - 'title': Title Slide (large title + subtitle)
                - 'content': Title and Content (title + bullet list)
                - 'two_content': Two Content (title + two columns)
                - 'blank': Blank slide
                - 'section': Section Header
                - 'title_only': Title Only
            title: Slide title text
            body: Main body text (use '\\n' for bullet points)
            subtitle: Subtitle text (for 'title' layout)

        Returns:
            SlideInfo with slide index and layout
        """
        prs = self._get_or_load_presentation(filename)

        layout_map = {
            "title": 0,
            "content": 1,
            "section": 2,
            "two_content": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6,
        }
        layout_idx = layout_map.get(layout, 1)

        # Clamp to available layouts
        layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        # Set body/content placeholder
        if body:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:  # body placeholder
                    placeholder.text = body
                    break

        # Set subtitle (for title layout, idx=1 is subtitle)
        if subtitle and layout == "title":
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:
                    placeholder.text = subtitle
                    break

        slide_index = len(prs.slides) - 1

        # Save
        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        logger.info("Slide added", filename=filename, slide_index=slide_index, layout=layout)
        return SlideInfo(slide_index=slide_index, layout=layout, title=title)

    @mcp_tool(description="Add a text box to a slide in a PowerPoint presentation")
    async def add_text_box(
        self,
        filename: str,
        slide_index: int,
        text: str,
        left_inches: float = 1.0,
        top_inches: float = 1.0,
        width_inches: float = 8.0,
        height_inches: float = 1.5,
        font_size_pt: int = 18,
        bold: bool = False,
        font_color_hex: str | None = None,
    ) -> str:
        """
        Add a text box to a specific slide.

        Args:
            filename: Presentation filename
            slide_index: Zero-based slide index
            text: Text content
            left_inches: Left position in inches
            top_inches: Top position in inches
            width_inches: Width in inches
            height_inches: Height in inches
            font_size_pt: Font size in points
            bold: Whether text is bold
            font_color_hex: Hex color string e.g. 'FF0000' for red

        Returns:
            Success message
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # Coerce numeric parameters (LLMs may pass them as strings)
        slide_index = self._coerce_int(slide_index, "slide_index")
        left_inches = self._coerce_float(left_inches, "left_inches")
        top_inches = self._coerce_float(top_inches, "top_inches")
        width_inches = self._coerce_float(width_inches, "width_inches")
        height_inches = self._coerce_float(height_inches, "height_inches")
        font_size_pt = self._coerce_int(font_size_pt, "font_size_pt")

        prs = self._get_or_load_presentation(filename)

        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range (total: {len(prs.slides)})")

        slide = prs.slides[slide_index]
        txBox = slide.shapes.add_textbox(
            Inches(left_inches), Inches(top_inches),
            Inches(width_inches), Inches(height_inches),
        )
        tf = txBox.text_frame
        tf.text = text

        para = tf.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold

        if font_color_hex:
            run.font.color.rgb = RGBColor.from_string(font_color_hex)

        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        return f"Text box added to slide {slide_index} in {filename}"

    @mcp_tool(description="Add a table to a slide in a PowerPoint presentation")
    async def add_table(
        self,
        filename: str,
        slide_index: int,
        headers: list[str],
        rows: list[list[str]],
        left_inches: float = 1.0,
        top_inches: float = 2.0,
        width_inches: float = 10.0,
        height_inches: float = 3.0,
    ) -> str:
        """
        Add a table to a specific slide.

        Args:
            filename: Presentation filename
            slide_index: Zero-based slide index
            headers: List of column header strings
            rows: List of rows, each row is a list of cell strings
            left_inches: Left position in inches
            top_inches: Top position in inches
            width_inches: Width in inches
            height_inches: Height in inches

        Returns:
            Success message
        """
        from pptx.util import Inches

        # Coerce numeric parameters (LLMs may pass them as strings)
        slide_index = self._coerce_int(slide_index, "slide_index")
        left_inches = self._coerce_float(left_inches, "left_inches")
        top_inches = self._coerce_float(top_inches, "top_inches")
        width_inches = self._coerce_float(width_inches, "width_inches")
        height_inches = self._coerce_float(height_inches, "height_inches")

        prs = self._get_or_load_presentation(filename)

        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")

        slide = prs.slides[slide_index]

        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left_inches), Inches(top_inches),
            Inches(width_inches), Inches(height_inches),
        ).table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header

        # Set data rows
        for row_idx, row in enumerate(rows):
            for col_idx, cell_value in enumerate(row):
                if col_idx < num_cols:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_value)

        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        return f"Table ({num_rows}x{num_cols}) added to slide {slide_index} in {filename}"

    @mcp_tool(description="Add an image to a slide in a PowerPoint presentation")
    async def add_image(
        self,
        filename: str,
        slide_index: int,
        image_path: str,
        left_inches: float = 1.0,
        top_inches: float = 1.5,
        width_inches: float | None = None,
        height_inches: float | None = None,
    ) -> str:
        """
        Add an image to a specific slide.

        Args:
            filename: Presentation filename
            slide_index: Zero-based slide index
            image_path: Path to image file (relative to workspace)
            left_inches: Left position in inches
            top_inches: Top position in inches
            width_inches: Width in inches (None = auto)
            height_inches: Height in inches (None = auto)

        Returns:
            Success message
        """
        from pptx.util import Inches

        # Coerce numeric parameters (LLMs may pass them as strings)
        slide_index = self._coerce_int(slide_index, "slide_index")
        left_inches = self._coerce_float(left_inches, "left_inches")
        top_inches = self._coerce_float(top_inches, "top_inches")
        if width_inches is not None:
            width_inches = self._coerce_float(width_inches, "width_inches")
        if height_inches is not None:
            height_inches = self._coerce_float(height_inches, "height_inches")

        prs = self._get_or_load_presentation(filename)

        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")

        img_path = self._resolve_path(image_path)
        if not img_path.exists():
            raise FileNotFoundError(f"Image not found: {image_path}")

        slide = prs.slides[slide_index]

        width = Inches(width_inches) if width_inches else None
        height = Inches(height_inches) if height_inches else None

        slide.shapes.add_picture(
            str(img_path),
            Inches(left_inches), Inches(top_inches),
            width=width, height=height,
        )

        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        return f"Image '{image_path}' added to slide {slide_index} in {filename}"

    @mcp_tool(description="Set the background color of a slide")
    async def set_slide_background(
        self,
        filename: str,
        slide_index: int,
        color_hex: str,
    ) -> str:
        """
        Set the background color of a slide.

        Args:
            filename: Presentation filename
            slide_index: Zero-based slide index
            color_hex: Hex color string e.g. '1F4E79' for dark blue

        Returns:
            Success message
        """
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from lxml import etree

        # Coerce numeric parameters (LLMs may pass them as strings)
        slide_index = self._coerce_int(slide_index, "slide_index")

        prs = self._get_or_load_presentation(filename)

        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")

        slide = prs.slides[slide_index]
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        return f"Background color set to #{color_hex} on slide {slide_index} in {filename}"

    @mcp_tool(description="Get information about a PowerPoint presentation")
    async def get_presentation_info(
        self,
        filename: str,
    ) -> dict[str, Any]:
        """
        Get information about an existing presentation.

        Args:
            filename: Presentation filename

        Returns:
            Dictionary with slide count, slide titles, etc.
        """
        prs = self._get_or_load_presentation(filename)

        slides_info = []
        for i, slide in enumerate(prs.slides):
            title = None
            if slide.shapes.title:
                title = slide.shapes.title.text
            slides_info.append({"index": i, "title": title})

        return {
            "filename": filename,
            "slide_count": len(prs.slides),
            "slides": slides_info,
            "width_inches": round(prs.slide_width.inches, 2),
            "height_inches": round(prs.slide_height.inches, 2),
        }

    @mcp_tool(description="Save a PowerPoint presentation to disk")
    async def save_presentation(
        self,
        filename: str,
        output_path: str | None = None,
    ) -> str:
        """
        Save the presentation to disk.

        Args:
            filename: Source presentation filename (in memory or on disk)
            output_path: Optional different output path

        Returns:
            Success message with saved path
        """
        prs = self._get_or_load_presentation(filename)
        save_path = self._resolve_path(output_path or filename)
        save_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(save_path))

        logger.info("Presentation saved", path=str(save_path))
        return f"Presentation saved to {save_path}"

    @mcp_tool(description="Create a complete PowerPoint presentation from a structured outline")
    async def create_presentation_from_outline(
        self,
        filename: str,
        title: str,
        subtitle: str | None = None,
        slides: list[dict[str, Any]] | None = None,
    ) -> PresentationInfo:
        """
        Create a complete presentation from a structured outline in one call.

        Args:
            filename: Output filename (e.g. 'report.pptx')
            title: Presentation title for the cover slide
            subtitle: Optional subtitle for the cover slide
            slides: List of slide dicts, each with:
                - 'title': Slide title (str)
                - 'layout': Layout name (str, default 'content')
                - 'body': Body text with newlines for bullets (str, optional)

        Returns:
            PresentationInfo

        Example slides:
            [
                {"title": "Agenda", "layout": "content", "body": "Point 1\\nPoint 2\\nPoint 3"},
                {"title": "Section 1", "layout": "section"},
                {"title": "Summary", "layout": "content", "body": "Key takeaway"},
            ]
        """
        info = await self.create_presentation(filename, title=None)

        pptx = self._get_pptx()
        prs = self._presentations[filename]

        # Add title/cover slide
        slide_layout = prs.slide_layouts[0]
        cover = prs.slides.add_slide(slide_layout)
        if cover.shapes.title:
            cover.shapes.title.text = title
        if subtitle:
            for ph in cover.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = subtitle
                    break

        slide_count = 1

        # Coerce slides: LLMs may pass the list as a JSON string, or individual
        # slide defs as JSON strings instead of dicts.
        if isinstance(slides, str):
            try:
                slides = _json.loads(slides)
            except _json.JSONDecodeError:
                raise ValueError(f"'slides' must be a list of dicts, got unparseable string: {slides!r}")

        for slide_def in (slides or []):
            # Handle individual slide defs that are JSON strings
            if isinstance(slide_def, str):
                try:
                    slide_def = _json.loads(slide_def)
                except _json.JSONDecodeError:
                    logger.warning("Skipping unparseable slide definition", raw=slide_def)
                    continue
            if not isinstance(slide_def, dict):
                logger.warning("Skipping non-dict slide definition", type=type(slide_def).__name__)
                continue

            layout = slide_def.get("layout", "content")
            s_title = slide_def.get("title")
            body = slide_def.get("body")

            await self.add_slide(filename, layout=layout, title=s_title, body=body)
            slide_count += 1

        file_path = self._resolve_path(filename)
        prs.save(str(file_path))

        logger.info("Presentation created from outline", filename=filename, slides=slide_count)
        return PresentationInfo(path=str(file_path), slide_count=slide_count, title=title)

    def get_tools(self) -> list[Callable]:
        """Get all tools as a list of callables."""
        return [
            self.create_presentation,
            self.add_slide,
            self.add_text_box,
            self.add_table,
            self.add_image,
            self.set_slide_background,
            self.get_presentation_info,
            self.save_presentation,
            self.create_presentation_from_outline,
        ]
