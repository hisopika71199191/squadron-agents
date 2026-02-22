from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define Modern/Energetic Colors
PRIMARY_COLOR = RGBColor(255, 20, 147)  # Deep Pink
SECONDARY_COLOR = RGBColor(0, 0, 0)     # Black
ACCENT_COLOR = RGBColor(255, 215, 0)    # Gold
BG_COLOR = RGBColor(255, 255, 255)      # White

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = 'Arial'
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'
    return slide

def add_content_slide(prs, title_text, content_points, image_placeholder_text="[Insert Product Image Here]"):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = 'Arial'
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = content_points[0]
    tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    tf.paragraphs[0].font.name = 'Arial'
    
    for point in content_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.color.rgb = SECONDARY_COLOR
        p.font.name = 'Arial'
        p.level = 0
        
    # Add a shape for image placeholder
    left = Inches(5)
    top = Inches(2)
    width = Inches(4)
    height = Inches(3)
    shape = slide.shapes.add_shape(
        1, left, top, width, height # msoShapeRectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = PRIMARY_COLOR
    shape.text_frame.text = image_placeholder_text
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].font.bold = True
    
    return slide

# Slide 1: Title
add_title_slide(
    prs, 
    "Chow Sang Sang 2026: Sparkle for the Modern Lady", 
    "Top 2 Product Recommendations for the Energetic Youth"
)

# Slide 2: Product 1 - Promessa (Hypothetical 2026 Edition based on trends)
product1_points = [
    "Collection: Promessa 'Eternal Bloom' 2026 Edition",
    "Why it fits 2026 Trends: Combines classic diamond solitaire with modern floral motifs.",
    "Key Feature: Interlocking rings symbolizing endless love, perfect for the independent young woman.",
    "Design Style: Minimalist yet statement-making, featuring ethically sourced diamonds."
]
add_content_slide(prs, "Recommendation 1: Promessa 'Eternal Bloom'", product1_points)

# Slide 3: Product 2 - Infini Love Diamond (Pink Twist)
product2_points = [
    "Collection: Infini Love Diamond 'Pink Twist' Series",
    "Why it fits 2026 Trends: Aligns with the 'Modern Pearls' and 'Statement Chokers' trend with a twist.",
    "Key Feature: Unique spiral setting enhancing brilliance, available in rose gold for a youthful glow.",
    "Design Style: Dynamic and fluid lines representing energy and movement."
]
add_content_slide(prs, "Recommendation 2: Infini Love 'Pink Twist'", product2_points)

# Save the presentation
output_file = "ChowSangSang_2026_Recommendations.pptx"
prs.save(output_file)
print(f"Presentation saved as {output_file}")