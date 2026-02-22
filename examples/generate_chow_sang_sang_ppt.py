from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define modern, energetic colors
primary_color = RGBColor(255, 20, 147)  # Deep Pink
secondary_color = RGBColor(255, 215, 0)  # Gold
bg_color = RGBColor(255, 255, 255)      # White
text_color = RGBColor(50, 50, 50)       # Dark Grey

# Helper to set slide background
def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Chow Sang Sang 2026\nTrend Recommendations"
subtitle.text = "Curated Selection for the Modern Young Lady\nEnergetic • Modern • Timeless"

# Style Title
tf = title.text_frame
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = primary_color
p.font.name = 'Arial'

# Style Subtitle
p_sub = subtitle.text_frame.paragraphs[0]
p_sub.font.size = Pt(24)
p_sub.font.color.rgb = text_color
p_sub.font.name = 'Arial'

# Slide 2: Product 1 - The 'Neon Galaxy' Collection (Hypothetical 2026 Line)
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Pick 1: 'Neon Galaxy' Diamond Pendant"
title.text_frame.paragraphs[0].font.color.rgb = primary_color
title.text_frame.paragraphs[0].font.bold = True

content_text = (
    "Why it fits 2026:\n"
    "• Design: Features a futuristic starburst motif with lab-grown diamonds, appealing to eco-conscious youth.\n"
    "• Style: Energetic geometric lines paired with 18K white gold.\n"
    "• Occasion: Perfect for layering or making a statement at evening events.\n"
    "• Target Vibe: Bold, independent, and sparkling."
)
content.text = content_text

# Slide 3: Product 2 - The 'Lucky Gallop' Heritage Series
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Pick 2: 'Lucky Gallop' Horse Charm Bracelet"
title.text_frame.paragraphs[0].font.color.rgb = primary_color
title.text_frame.paragraphs[0].font.bold = True

content_text = (
    "Why it fits 2026 (Year of the Horse):\n"
    "• Design: A minimalist, abstract horse silhouette in rose gold, modernizing the traditional zodiac symbol.\n"
    "• Style: Delicate chain with a vibrant enamel accent for an energetic pop of color.\n"
    "• Meaning: Symbolizes speed, success, and freedom for the career-driven young lady.\n"
    "• Target Vibe: Chic, meaningful, and versatile."
)
content.text = content_text

# Save
prs.save('Chow_Sang_Sang_2026_Recommendation.pptx')
print("Presentation generated successfully.")